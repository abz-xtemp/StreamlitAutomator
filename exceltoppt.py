import pandas as pd
import win32com.client
import pythoncom
from pptx import Presentation
from pptx.util import Inches
import openpyxl
from PIL import Image
import tempfile
import os
import sys
import logging
import time

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def process(ppt_file, excel_file, sheet_name, cell_range, slide_number, slide_width, slide_height, left, top, excel_password=None):
    """
    Process Excel data and insert it into a PowerPoint presentation.
    
    Parameters:
    - ppt_file: PowerPoint file object or path (can be None to create new presentation)
    - excel_file: Excel file object or path
    - sheet_name: Name of the Excel sheet
    - cell_range: Excel range to copy (e.g., "A1:H10")
    - slide_number: Slide number to paste into (1-based index)
    - slide_width: Width of the pasted image in inches
    - slide_height: Height of the pasted image in inches
    - left: Left position in inches
    - top: Top position in inches
    - excel_password: Optional password for protected Excel files
    
    Returns:
    - Path to the saved PowerPoint file or error message
    """
    temp_files = []  # Keep track of temp files to clean up later
    xlApp = None
    
    try:
        # Initialize COM for this thread
        pythoncom.CoInitialize()
        
        # Validate inputs
        if excel_file is None:
            raise ValueError("Excel file is missing")
            
        # Save uploaded files as temporary files
        if hasattr(ppt_file, 'read'):
            # It's a file-like object
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_ppt:
                temp_ppt.write(ppt_file.read())
                temp_ppt_path = temp_ppt.name
                temp_files.append(temp_ppt_path)
        elif ppt_file is not None and isinstance(ppt_file, str):
            # It's a file path
            temp_ppt_path = ppt_file
        else:
            # Create a new presentation
            temp_ppt_path = None
            
        if hasattr(excel_file, 'read'):
            # It's a file-like object
            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as temp_xlsx:
                temp_xlsx.write(excel_file.read())
                temp_xlsx_path = temp_xlsx.name
                temp_files.append(temp_xlsx_path)
        elif isinstance(excel_file, str):
            # It's a file path
            temp_xlsx_path = excel_file
        else:
            raise ValueError("Excel file must be a file-like object or a file path")
            
        # Validate the Excel file exists
        if not os.path.exists(temp_xlsx_path):
            raise FileNotFoundError(f"Excel file not found at path: {temp_xlsx_path}")
            
        # Load Excel
        logger.info(f"Opening Excel file: {temp_xlsx_path}")
        xlApp = win32com.client.Dispatch("Excel.Application")
        xlApp.Visible = True  # Set to True for debugging
        xlApp.DisplayAlerts = False  # Suppress alerts
        
        if excel_password:
            wb = xlApp.Workbooks.Open(temp_xlsx_path, False, True, None, excel_password)
        else:
            wb = xlApp.Workbooks.Open(temp_xlsx_path, False, True)
            
        # Validate sheet name exists
        sheet_names = [sheet.Name for sheet in wb.Sheets]
        if sheet_name not in sheet_names:
            available_sheets = ", ".join(sheet_names)
            raise ValueError(f"Sheet '{sheet_name}' not found. Available sheets: {available_sheets}")
            
        ws = wb.Sheets(sheet_name)
        
        # Ensure the sheet is active and visible
        ws.Activate()
        
        # Check if the range exists and contains data
        try:
            range_obj = ws.Range(cell_range)
            if range_obj.Count == 0:
                raise ValueError(f"The specified range '{cell_range}' is empty")
        except Exception as e:
            raise ValueError(f"Problem with cell range '{cell_range}': {str(e)}")
        
        logger.info(f"Copying range {cell_range} from sheet {sheet_name}")
        
        # Try to export the range to an image using a chart object
        try:
            logger.info("Exporting range as image using chart object")
            chart_obj = ws.ChartObjects().Add(10, 10, 300, 300)  # Add a chart object temporarily
            chart = chart_obj.Chart
            chart.SetSourceData(range_obj)
            
            # Export the chart as an image
            temp_img_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
            chart.Export(temp_img_path)
            chart_obj.Delete()  # Clean up the chart object
            
            # Load the exported image
            img = Image.open(temp_img_path)
            logger.info(f"Successfully created image: {temp_img_path}")
        except Exception as e:
            logger.error(f"Failed to export range as image: {str(e)}")
            raise ValueError("Failed to export range as image from Excel")
        
        # Save the image temporarily
        temp_files.append(temp_img_path)
        
        # Load or Create PowerPoint
        logger.info("Working with PowerPoint presentation")
        if temp_ppt_path:
            prs = Presentation(temp_ppt_path)
        else:
            prs = Presentation()
            
        # Validate slide number
        if slide_number < 1:
            raise ValueError("Slide number must be at least 1")
            
        # Adjust slide count if needed
        while len(prs.slides) < slide_number:
            prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            
        slide = prs.slides[slide_number - 1]  # Get target slide
        
        # Insert Image
        logger.info(f"Adding image to slide {slide_number} at position ({left}, {top}) with dimensions {slide_width}x{slide_height}")
        slide.shapes.add_picture(temp_img_path, Inches(left), Inches(top), Inches(slide_width), Inches(slide_height))
        
        # Save Presentation
        output_ppt = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx").name
        temp_files.append(output_ppt)
        prs.save(output_ppt)
        logger.info(f"Saved presentation to {output_ppt}")
        
        # Close Excel
        wb.Close(SaveChanges=False)
        
        return output_ppt
        
    except Exception as e:
        logger.error(f"Error in process function: {str(e)}", exc_info=True)
        return f"Error: {str(e)}"
        
    finally:
        # Clean up Excel application
        if xlApp:
            try:
                xlApp.Quit()
            except:
                pass
                
        # Uninitialize COM
        pythoncom.CoUninitialize()
                
        # Optionally clean up temporary files (uncomment if you want to delete them)
        # for temp_file in temp_files:
        #     try:
        #         if os.path.exists(temp_file):
        #             os.remove(temp_file)
        #     except:
        #         pass


def main():
    """
    Command-line interface for the process function.
    
    Usage:
    python script.py excel.xlsx PowerPoint.pptx "Sheet1" "A1:H10" 1 5 3 1 1
    """
    if len(sys.argv) < 9:
        print("Usage: python script.py excel.xlsx [powerpoint.pptx] sheet_name cell_range slide_number width height left top [excel_password]")
        return
        
    excel_path = sys.argv[1]
    
    if len(sys.argv) >= 10:
        ppt_path = sys.argv[2]
        start_idx = 3
    else:
        ppt_path = None
        start_idx = 2
        
    sheet_name = sys.argv[start_idx]
    cell_range = sys.argv[start_idx + 1]
    slide_number = int(sys.argv[start_idx + 2])
    width = float(sys.argv[start_idx + 3])
    height = float(sys.argv[start_idx + 4])
    left = float(sys.argv[start_idx + 5])
    top = float(sys.argv[start_idx + 6])
    
    excel_password = sys.argv[start_idx + 7] if len(sys.argv) > start_idx + 7 else None
    
    result = process(ppt_path, excel_path, sheet_name, cell_range, slide_number, width, height, left, top, excel_password)
    print(result)


if __name__ == "__main__":
    # Initialize COM at the application level
    pythoncom.CoInitialize()
    try:
        main()
    finally:
        # Uninitialize COM when done
        pythoncom.CoUninitialize()
